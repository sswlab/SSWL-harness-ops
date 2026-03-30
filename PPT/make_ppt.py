#!/usr/bin/env python3
"""SSWL AI 하네스 v2.0 소개 팜플렛 — PPT 생성."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x28, 0x35, 0x93)
INDIGO = RGBColor(0x39, 0x49, 0xAB)
LIGHT_BG = RGBColor(0xE8, 0xEA, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
TEAL = RGBColor(0x00, 0x69, 0x5C)
PURPLE = RGBColor(0x4A, 0x14, 0x8C)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_TEAL = RGBColor(0xE0, 0xF2, 0xF1)
LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x1A, 0x23, 0x7E)
TABLE_ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_BORDER = RGBColor(0xC5, 0xCA, 0xE9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT_NAME = "Noto Sans CJK KR"


# ── 헬퍼 함수 ──
def set_font(run, size=14, bold=False, color=BLACK, name=FONT_NAME, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    run.font.italic = italic


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table


def style_table(table, headers, data, col_widths=None):
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        set_font(run, size=12, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j > 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            set_font(run, size=11, color=DARK_GRAY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT if i % 2 == 1 else WHITE
    tblPr = table._tbl.tblPr
    if tblPr is None:
        tblPr = table._tbl._new_tblPr()
    tblPr.set("firstRow", "1")
    tblPr.set("bandRow", "1")


def add_paragraph(tf, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, space_after=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return run


def add_step_boxes(slide, steps, start_x, start_y, accent_color, dark_color):
    for i, (title, desc) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = start_x + col * 1.9
        y = start_y + row * 2.1
        add_rect(slide, x, y, 1.7, 1.8, WHITE, accent_color)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.65), Inches(y + 0.1), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_color
        circle.line.fill.background()
        cp = circle.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        set_font(cr, size=13, bold=True, color=WHITE)
        tb2 = add_textbox(slide, x + 0.1, y + 0.55, 1.5, 0.4)
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = title
        set_font(r2, size=12, bold=True, color=dark_color)
        tb3 = add_textbox(slide, x + 0.1, y + 0.95, 1.5, 0.8)
        tb3.text_frame.word_wrap = True
        p3 = tb3.text_frame.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = desc
        set_font(r3, size=10, color=GRAY)


# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.333, 0.08, INDIGO)

tb = add_textbox(slide, 1, 1.5, 11.333, 0.6)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "태양 및 우주환경연구실을 위한 연구 지원 AI 시스템"
set_font(run, size=18, color=RGBColor(0xC5, 0xCA, 0xE9))

tb = add_textbox(slide, 1, 2.2, 11.333, 1.5)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "SSWL AI Harness v2.0"
set_font(run, size=48, bold=True, color=WHITE)

add_rect(slide, 5, 3.7, 3.333, 0.04, INDIGO)

tb = add_textbox(slide, 1, 4.0, 11.333, 0.6)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "모델 아카이빙  ·  연구 업무 수행  ·  아이디어→실험 결과  ·  논문 초안"
set_font(run, size=20, color=RGBColor(0x9F, 0xA8, 0xDA))

tb = add_textbox(slide, 2, 5.0, 9.333, 0.8)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '"연구실에서 개발된 모델을 누구나 활용하고, 아이디어를 빠르게 검증합니다."'
set_font(run, size=15, color=RGBColor(0x9F, 0xA8, 0xDA), italic=True)

tb = add_textbox(slide, 1, 6.5, 11.333, 0.5)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "2026  ·  태양 및 우주환경연구실"
set_font(run, size=13, color=RGBColor(0x79, 0x86, 0xCB))


# ============================================================
# 슬라이드 2: 왜 필요한가?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "이 시스템은 왜 필요한가요?"
set_font(run, size=30, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 2.5, 0.05, INDIGO)

headers = ["기존 방식", "문제점"]
data = [
    ["연구실 모델이 개발자 PC에만 존재", "개발자 부재 시 아무도 실행 불가"],
    ["모델 실행 방법이 구전으로 전수", "새 연구자 합류 시 재교육 필요"],
    ["데이터 수집을 사이트별로 수동", "연구 시간 상당 부분이 데이터 수집에 소요"],
    ["아이디어에서 실험까지 오래 걸림", "아이디어 → 데이터 → 실행 → 분석이 수일~수주"],
    ["모델 버전 관리 부재", "어떤 버전으로 어떤 결과를 냈는지 추적 불가"],
]
tbl = add_table(slide, 6, 2, 0.8, 1.5, 7.5, 3.0)
style_table(tbl, headers, data, col_widths=[3.5, 4.0])

box = add_rect(slide, 8.8, 1.5, 3.8, 3.0, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 9.0, 1.7, 3.4, 2.6)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "SSWL AI 하네스 v2.0은"
set_font(run, size=14, bold=True, color=NAVY)
add_paragraph(tf, "", size=6)
add_paragraph(tf, "모델을 표준화·아카이빙하고", size=13, color=DARK_GRAY)
add_paragraph(tf, "누구나 자연어로 실행하며", size=13, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "아이디어를 빠르게 실험 결과로", size=13, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "전환합니다.", size=13, color=DARK_GRAY)

box = add_rect(slide, 0.8, 5.0, 11.7, 0.9, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 1.0, 5.1, 11.3, 0.7)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "모델은 자산화  ·  누구나 활용  ·  아이디어에서 결과까지 당일 확인"
set_font(run, size=18, bold=True, color=INDIGO)


# ============================================================
# 슬라이드 3: 시스템 한눈에 보기 — 3가지 모드
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "시스템 한눈에 보기 — 세 가지 핵심 기능"
set_font(run, size=30, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 2.5, 0.05, INDIGO)

# 기능 1: 연구 업무 수행
add_rect(slide, 0.5, 1.5, 3.9, 5.2, LIGHT_BG, RGBColor(0x90, 0xCA, 0xF9))
tb = add_textbox(slide, 0.6, 1.6, 3.7, 0.5)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "① 연구 업무 수행"
set_font(run, size=16, bold=True, color=BLUE)

steps = [
    ("연구자 요청", "자연어로\n작업 요청"),
    ("계획 수립", "AI가 실행계획\n자동 생성"),
    ("승인", "연구자가\n계획 검토·승인"),
    ("자동 실행", "데이터 수집 →\n모델 실행"),
    ("결과 보고", "보고서 생성\n시각화 포함"),
    ("후속 제안", "다음 분석\n자동 제안"),
]
for i, (title, desc) in enumerate(steps):
    row = i // 2
    col = i % 2
    x = 0.6 + col * 1.9
    y = 2.2 + row * 1.55
    add_rect(slide, x, y, 1.7, 1.35, WHITE, INDIGO)
    tb2 = add_textbox(slide, x + 0.1, y + 0.1, 1.5, 0.35)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"{i+1}. {title}"
    set_font(r2, size=11, bold=True, color=NAVY)
    tb3 = add_textbox(slide, x + 0.1, y + 0.5, 1.5, 0.7)
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = desc
    set_font(r3, size=10, color=GRAY)

# 기능 2: 모델 아카이빙
add_rect(slide, 4.7, 1.5, 3.9, 5.2, LIGHT_GREEN, RGBColor(0xA5, 0xD6, 0xA7))
tb = add_textbox(slide, 4.8, 1.6, 3.7, 0.5)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "② 모델 아카이빙 & 활용"
set_font(run, size=16, bold=True, color=GREEN)

archive_items = [
    ("모델 등록", "코드 분석 →\n모델 카드 자동 생성"),
    ("모델 목록", "등록된 모델\n표/카드로 조회"),
    ("모델 활용", "자연어 요청 →\n자동 실행·결과 보고"),
    ("버전 관리", "업데이트 시\n이전 버전 보존"),
]
for i, (title, desc) in enumerate(archive_items):
    row = i // 2
    col = i % 2
    x = 4.8 + col * 1.9
    y = 2.2 + row * 1.55
    add_rect(slide, x, y, 1.7, 1.35, WHITE, GREEN)
    tb2 = add_textbox(slide, x + 0.1, y + 0.1, 1.5, 0.35)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    set_font(r2, size=11, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
    tb3 = add_textbox(slide, x + 0.1, y + 0.5, 1.5, 0.7)
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = desc
    set_font(r3, size=10, color=GRAY)

# 핵심 문구
tb = add_textbox(slide, 4.8, 5.65, 3.7, 0.9)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "개발에 참여하지 않은\n연구자도 동일하게 사용 가능"
set_font(run, size=12, bold=True, color=GREEN)

# 기능 3: 아이디어 → 실험
add_rect(slide, 8.9, 1.5, 3.9, 5.2, LIGHT_PURPLE, RGBColor(0xCE, 0x93, 0xD8))
tb = add_textbox(slide, 9.0, 1.6, 3.7, 0.5)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "③ 아이디어 → 실험 결과"
set_font(run, size=16, bold=True, color=PURPLE)

idea_items = [
    ("아이디어 분석", "핵심 질문 추출\n가설 수립"),
    ("실험 설계", "Baseline vs\nExperiment 설계"),
    ("실험 실행", "데이터 수집 →\n모델 실행 → 비교"),
    ("결과 + 논문", "통계 분석 보고\n논문 초안 작성"),
]
for i, (title, desc) in enumerate(idea_items):
    row = i // 2
    col = i % 2
    x = 9.0 + col * 1.9
    y = 2.2 + row * 1.55
    add_rect(slide, x, y, 1.7, 1.35, WHITE, PURPLE)
    tb2 = add_textbox(slide, x + 0.1, y + 0.1, 1.5, 0.35)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    set_font(r2, size=11, bold=True, color=PURPLE)
    tb3 = add_textbox(slide, x + 0.1, y + 0.5, 1.5, 0.7)
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = desc
    set_font(r3, size=10, color=GRAY)

# 핵심 문구
tb = add_textbox(slide, 9.0, 5.65, 3.7, 0.9)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '"교수님 아이디어 말씀하시면\n간단한 실험 결과 보여드립니다"'
set_font(run, size=12, bold=True, color=PURPLE, italic=True)


# ============================================================
# 슬라이드 4: 연구 업무 수행 상세
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = '기능 1. 연구 업무 수행 — "말하면 해줍니다"'
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 좌측: 요청 예시
add_rect(slide, 0.8, 1.4, 5.5, 2.5, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 1.0, 1.5, 5.1, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "요청 예시"
set_font(run, size=16, bold=True, color=NAVY)

examples = [
    '"최근 1달 HMI 데이터로 synoptic map 만들어줘"',
    '"AIA 193 이미지에서 코로나홀 찾아줘"',
    '"이걸로 PFSS 시뮬레이션 돌려줘"',
]
tb = add_textbox(slide, 1.0, 2.0, 5.1, 1.8)
tf = tb.text_frame
tf.word_wrap = True
for ex in examples:
    add_paragraph(tf, ex, size=13, color=DARK_GRAY, space_after=8)

# 우측: 처리 흐름
add_rect(slide, 6.8, 1.4, 5.8, 2.5, WHITE, TABLE_BORDER)
tb = add_textbox(slide, 7.0, 1.5, 5.4, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "처리 흐름"
set_font(run, size=16, bold=True, color=NAVY)

flow_steps = [
    "요청 → query-planner가 실행 계획 수립",
    "연구자 승인 (반드시 승인 후 실행)",
    "task-executor가 자동 실행 (데이터→모델→후처리)",
    "result-reporter가 결과 보고 + 후속 제안",
]
tb = add_textbox(slide, 7.0, 2.0, 5.4, 1.8)
tf = tb.text_frame
tf.word_wrap = True
for i, step in enumerate(flow_steps):
    add_paragraph(tf, f"{i+1}. {step}", size=12, color=DARK_GRAY, space_after=6)

# 하단: 시나리오 예시
add_rect(slide, 0.8, 4.2, 11.7, 2.5, RGBColor(0xFD, 0xFD, 0xE8), RGBColor(0xE0, 0xE0, 0xE0))
tb = add_textbox(slide, 1.0, 4.3, 11.3, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "시나리오: Synoptic Map → PFSS 연속 작업"
set_font(run, size=14, bold=True, color=NAVY)

tb = add_textbox(slide, 1.0, 4.7, 11.3, 1.8)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '연구자: "최근 1달 SO/PHI와 SDO/HMI 데이터로 synoptic map 만들어줘"', size=11, color=DARK_GRAY, space_after=4)
add_paragraph(tf, '→ AI: 계획 제시 (4단계, 40~60분 예상) → 승인 후 자동 실행 → 결과 보고', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '→ "이 map으로 PFSS 시뮬레이션 가능합니다"', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '연구자: "이걸로 PFSS 돌려줘" → 이전 결과 자동 참조하여 즉시 실행', size=11, color=DARK_GRAY, space_after=4)


# ============================================================
# 슬라이드 5: 모델 아카이빙 상세
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = '기능 2. 모델 아카이빙 — "개발하지 않았어도 사용합니다"'
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 좌측: 등록 흐름
add_rect(slide, 0.8, 1.4, 5.8, 2.2, LIGHT_GREEN, RGBColor(0xA5, 0xD6, 0xA7))
tb = add_textbox(slide, 1.0, 1.5, 5.4, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "모델 등록 흐름"
set_font(run, size=16, bold=True, color=GREEN)

tb = add_textbox(slide, 1.0, 2.0, 5.4, 1.5)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '개발자: "코로나홀 탐지 모델 등록해줘. 코드는 여기에."', size=12, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '→ AI가 코드 분석 (입출력, 의존성, 실행 방법 파악)', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '→ 모델 카드(사용 설명서) + 실행 래퍼 자동 생성', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '→ model_registry에 등록 완료 + 테스트 실행', size=11, color=GRAY)

# 우측: 활용 흐름
add_rect(slide, 7.0, 1.4, 5.8, 2.2, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 7.2, 1.5, 5.4, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "모델 활용 (개발 미참여 연구자)"
set_font(run, size=16, bold=True, color=NAVY)

tb = add_textbox(slide, 7.2, 2.0, 5.4, 1.5)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '연구자B: "코로나홀 탐지 모델로 어제 AIA 193 분석해줘"', size=12, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '→ 모델 카드 참조 → 필요 데이터 자동 수집 (JSOC)', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '→ 모델 자동 실행', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '→ 결과 시각화 및 보고', size=11, color=GRAY)

# 하단: 모델 목록 테이블
tb = add_textbox(slide, 0.8, 3.9, 5, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "등록 모델 목록 예시"
set_font(run, size=16, bold=True, color=NAVY)

headers = ["#", "모델명", "용도", "입력", "개발자"]
data = [
    ["1", "coronal_hole_detect", "코로나홀 탐지", "AIA 193 A", "김연구"],
    ["2", "dem_model", "DEM 분석", "AIA 6채널", "이연구"],
    ["3", "pix2pixcc_fsi", "EUI 채널 합성", "FSI 174+304", "윤준무"],
    ["4", "synoptic_map", "Synoptic map", "HMI 자기장", "(내장)"],
    ["5", "pfss_sim", "PFSS 시뮬레이션", "Synoptic map", "(내장)"],
]
tbl = add_table(slide, 6, 5, 0.8, 4.3, 11.7, 2.5)
style_table(tbl, headers, data, col_widths=[0.5, 2.5, 2.5, 2.5, 1.5])


# ============================================================
# 슬라이드 6: 아이디어 → 실험 결과 상세
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = '기능 3. 아이디어 → 실험 결과 — "말하면 실험합니다"'
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 시나리오 예시
add_rect(slide, 0.8, 1.4, 11.7, 5.3, RGBColor(0xFD, 0xFD, 0xFD), RGBColor(0xE0, 0xE0, 0xE0))

tb = add_textbox(slide, 1.0, 1.5, 11.3, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "시나리오: EUI DEM 연구에서 Magnetogram 추가 효과 검증"
set_font(run, size=15, bold=True, color=PURPLE)

# 좌측: 대화 흐름
tb = add_textbox(slide, 1.0, 2.0, 5.5, 4.5)
tf = tb.text_frame
tf.word_wrap = True

add_paragraph(tf, '교수님:', size=12, bold=True, color=PURPLE, space_after=2)
add_paragraph(tf, '"Solar Orbiter DEM 연구에서 Magnetogram을\n DL 입력에 추가하면 합성 채널이 좋아질까?"', size=11, color=DARK_GRAY, space_after=8)

add_paragraph(tf, 'AI → 아이디어 분석:', size=12, bold=True, color=NAVY, space_after=2)
add_paragraph(tf, '배경: EUI/FSI 2채널 → Pix2PixCC → 합성 5채널 → DEM', size=10, color=GRAY, space_after=2)
add_paragraph(tf, '제안: DL 입력에 magnetogram 추가 시 품질 개선?', size=10, color=GRAY, space_after=8)

add_paragraph(tf, 'AI → 실험 설계:', size=12, bold=True, color=NAVY, space_after=2)
add_paragraph(tf, 'Baseline: AIA 171+304 → Pix2PixCC → DEM', size=10, color=GRAY, space_after=2)
add_paragraph(tf, 'Experiment: 171+304+HMI mag → Modified Pix2PixCC → DEM', size=10, color=GRAY, space_after=2)
add_paragraph(tf, 'Reference: 실제 AIA 6채널 DEM (ground truth)', size=10, color=GRAY, space_after=8)

add_paragraph(tf, '교수님: "비교할 때 활동영역과 조용한 영역 따로 봐줘"', size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, '→ AI: 영역 분류 단계 추가하여 계획 수정', size=10, color=GRAY)

# 우측: 결과
add_rect(slide, 6.8, 2.0, 5.5, 2.5, LIGHT_PURPLE, RGBColor(0xCE, 0x93, 0xD8))
tb = add_textbox(slide, 7.0, 2.1, 5.1, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "실험 결과 (예시)"
set_font(run, size=14, bold=True, color=PURPLE)

tb = add_textbox(slide, 7.0, 2.5, 5.1, 1.8)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '합성 94 A 채널 CC: 0.87 → 0.91 개선', size=11, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '활동영역 DEM:', size=11, bold=True, color=PURPLE, space_after=2)
add_paragraph(tf, '  logT 6.0~6.5 EM 오차 15% 감소', size=11, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '조용한 영역:', size=11, bold=True, color=PURPLE, space_after=2)
add_paragraph(tf, '  유의미한 차이 없음', size=11, color=DARK_GRAY)

# 논문 초안 박스
add_rect(slide, 6.8, 4.8, 5.5, 1.7, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 7.0, 4.9, 5.1, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = '"이걸로 논문 초안 써줘"'
set_font(run, size=14, bold=True, color=NAVY)

tb = add_textbox(slide, 7.0, 5.3, 5.1, 1.1)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '→ Introduction: EUI DEM 한계 + DL 접근 배경', size=10, color=GRAY, space_after=2)
add_paragraph(tf, '→ Method: Pix2PixCC 수정, 학습 데이터, DEM 방법론', size=10, color=GRAY, space_after=2)
add_paragraph(tf, '→ Results: 채널 품질 비교 + DEM 비교 (영역별)', size=10, color=GRAY, space_after=2)
add_paragraph(tf, '→ Discussion: 자기장 정보 기여 해석 + 한계', size=10, color=GRAY)


# ============================================================
# 슬라이드 7: 연구실 논문 아카이브
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "연구실 논문 아카이브 (51편, 2019~2026)"
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

headers = ["연구 분야", "주요 내용", "논문 수"]
data = [
    ["딥러닝 이미지 변환/합성", "EUV↔자기장, 채널 합성, 백색광, He I", "8편"],
    ["태양 뒷면 자기장", "STEREO/EUV → farside magnetogram", "4편"],
    ["DEM 연구", "EUI 합성 채널 → DEM, Pixel-to-pixel", "2편"],
    ["우주기상 예보", "태양풍, IMF Bz, 플레어, 전리층 예측", "12편"],
    ["태양 자기장/코로나", "NLFFF, PFSS, Synoptic map, 3D 복원", "8편"],
    ["코로나 진동/파동", "플룸, 루프, Alfven 속도", "5편"],
    ["CME/SEP", "CME 질량, ICME, SEP 소스", "5편"],
    ["미션/관측 시스템", "L4 미션, 망원경 제어, 채널 최적화", "3편"],
]
tbl = add_table(slide, 9, 3, 0.8, 1.4, 7.0, 4.0)
style_table(tbl, headers, data, col_widths=[2.5, 3.5, 1.0])

# 우측: 주요 저널
add_rect(slide, 8.3, 1.4, 4.5, 2.5, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 8.5, 1.5, 4.1, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "게재 저널"
set_font(run, size=16, bold=True, color=NAVY)

tb = add_textbox(slide, 8.5, 1.9, 4.1, 1.8)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "Nature Astronomy — 2편", size=12, bold=True, color=DARK_GRAY, space_after=4)
add_paragraph(tf, "ApJ / ApJS / ApJL — 다수", size=12, color=DARK_GRAY, space_after=4)
add_paragraph(tf, "A&A, Space Weather", size=12, color=DARK_GRAY, space_after=4)
add_paragraph(tf, "Solar Physics, Remote Sensing", size=12, color=DARK_GRAY, space_after=4)
add_paragraph(tf, "JKAS, JASTP, PASP", size=12, color=DARK_GRAY)

# 활용 예시
add_rect(slide, 8.3, 4.2, 4.5, 2.3, RGBColor(0xFD, 0xFD, 0xE8), RGBColor(0xE0, 0xE0, 0xE0))
tb = add_textbox(slide, 8.5, 4.3, 4.1, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "AI에게 물어보세요"
set_font(run, size=14, bold=True, color=NAVY)

tb = add_textbox(slide, 8.5, 4.7, 4.1, 1.6)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '"플레어 예보 관련 논문 뭐가 있어?"', size=11, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '"태양 뒷면 자기장 연구 히스토리\n 정리해줘"', size=11, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '"DEM 관련 우리 논문 참고문헌에\n 넣으려는데 찾아줘"', size=11, color=DARK_GRAY)


# ============================================================
# 슬라이드 8: 기대 효과
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "기대 효과"
set_font(run, size=30, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 1.5, 0.05, INDIGO)

# 연구자 개인
headers = ["항목", "기존", "도입 후"]
data = [
    ["모델 사용", "개발자에게 의존, 구전 교육", "자연어로 누구나 즉시 실행"],
    ["데이터 수집", "여러 사이트 수동 다운로드", "AI가 자동 수집·검증"],
    ["아이디어 검증", "데이터 수집~분석까지 수일", "아이디어 제시 후 당일 결과"],
    ["모델 실행", "명령어 직접 입력, 환경 세팅", "자연어 요청으로 자동 실행"],
    ["논문 초안", "백지부터 직접 작성", "실험 결과 기반 골격 자동 생성"],
]
tbl = add_table(slide, 6, 3, 0.8, 1.3, 7.5, 2.8)
style_table(tbl, headers, data, col_widths=[1.5, 3.0, 3.0])

# 연구실 차원
add_rect(slide, 8.8, 1.3, 4.0, 2.8, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 9.0, 1.4, 3.6, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "연구실 차원"
set_font(run, size=16, bold=True, color=NAVY)

tb = add_textbox(slide, 9.0, 1.8, 3.6, 2.1)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "모델 자산화", size=13, bold=True, color=NAVY, space_after=2)
add_paragraph(tf, "개인 PC에 묻히지 않고 연구실 자산으로", size=10, color=GRAY, space_after=8)
add_paragraph(tf, "지식 전수", size=13, bold=True, color=NAVY, space_after=2)
add_paragraph(tf, "신규 연구자도 등록 모델 즉시 활용", size=10, color=GRAY, space_after=8)
add_paragraph(tf, "연구 가속", size=13, bold=True, color=NAVY, space_after=2)
add_paragraph(tf, "아이디어 → 초기 결과까지 시간 대폭 단축", size=10, color=GRAY, space_after=8)
add_paragraph(tf, "재현성 확보", size=13, bold=True, color=NAVY, space_after=2)
add_paragraph(tf, "모델 버전 관리로 결과 추적 가능", size=10, color=GRAY)

# 하단: v1 vs v2 비교
tb = add_textbox(slide, 0.8, 4.4, 5, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "v1.0 → v2.0 변경사항"
set_font(run, size=16, bold=True, color=NAVY)

headers = ["항목", "v1.0", "v2.0"]
data = [
    ["상시운용 (24h 감시)", "있음", "제거"],
    ["연구 업무 수행", "보조 모드", "핵심 모드"],
    ["모델 아카이빙", "없음", "신규 추가"],
    ["아이디어→실험", "없음", "신규 추가"],
    ["논문 초안 작성", "없음", "신규 추가"],
]
tbl = add_table(slide, 6, 3, 0.8, 4.8, 11.7, 2.3)
style_table(tbl, headers, data, col_widths=[3.5, 3.0, 3.0])


# ============================================================
# 슬라이드 9: 현재 한계 — 툴체인 실행 레이어 부재
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "현재 한계: 툴체인 실행 레이어 부재"
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 좌측: 있는 것 vs 없는 것
add_rect(slide, 0.5, 1.4, 6.0, 2.8, RGBColor(0xFF, 0xEB, 0xEE), RGBColor(0xEF, 0x9A, 0x9A))
tb = add_textbox(slide, 0.7, 1.5, 5.6, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "현재 구조의 문제"
set_font(run, size=16, bold=True, color=RGBColor(0xC6, 0x28, 0x28))

headers = ["있는 것 (설명서)", "없는 것 (실행 레이어)"]
data = [
    ['"data-fetcher는 JSOC에서 수집한다"', "JSOC 다운로드를 실행하는 도구"],
    ['"model-runner는 모델을 실행한다"', "python3 run.py를 자동 체이닝하는 로직"],
    ['"task-executor가 순차 실행한다"', "Step1→Step2 출력을 자동 연결하는 파이프라인"],
]
tbl = add_table(slide, 4, 2, 0.5, 2.1, 6.0, 1.8)
style_table(tbl, headers, data, col_widths=[3.0, 3.0])

# 우측: 결과
add_rect(slide, 6.8, 1.4, 6.0, 2.8, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 7.0, 1.5, 5.6, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "그래서 어떻게 되는가"
set_font(run, size=16, bold=True, color=NAVY)

tb = add_textbox(slide, 7.0, 2.0, 5.6, 2.0)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, 'AI가 매번 즉흥적으로 코드를 작성하거나 검색', size=12, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '같은 작업을 할 때마다 다른 방법으로 실행', size=12, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '연구실의 검증된 코드를 정확히 호출하지 못함', size=12, color=DARK_GRAY, space_after=6)
add_paragraph(tf, 'Step간 데이터 전달을 AI가 매번 추론해야 함', size=12, color=DARK_GRAY, space_after=6)
add_paragraph(tf, '→ 재현성 없음, 오류 빈발, 신뢰도 낮음', size=13, bold=True, color=RGBColor(0xC6, 0x28, 0x28))

# 하단: 어시웍스 참고
add_rect(slide, 0.5, 4.5, 12.3, 2.3, RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0x90, 0xCA, 0xF9))
tb = add_textbox(slide, 0.7, 4.6, 5.0, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "참고: AssiWorks의 4계층 모델"
set_font(run, size=15, bold=True, color=BLUE)

# 4계층 박스
layers_assi = [
    ("Tool", "실행 가능한 최소 단위\n(API, 스크립트, LLM)"),
    ("Flow", "Tool들을 순차/조건/\n병렬로 연결한 체인"),
    ("Agent", "자연어 이해 →\n적절한 Tool/Flow 선택"),
    ("Team", "여러 Agent가\n협업하여 복합 처리"),
]
for i, (name, desc) in enumerate(layers_assi):
    x = 0.7 + i * 3.05
    add_rect(slide, x, 5.1, 2.8, 1.5, WHITE, BLUE)
    tb2 = add_textbox(slide, x + 0.1, 5.15, 2.6, 0.35)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"Layer {i+1}: {name}"
    set_font(r2, size=12, bold=True, color=BLUE)
    tb3 = add_textbox(slide, x + 0.1, 5.5, 2.6, 0.9)
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = desc
    set_font(r3, size=10, color=GRAY)

# 화살표
for i in range(3):
    x = 3.35 + i * 3.05
    tb5 = add_textbox(slide, x, 5.55, 0.4, 0.4)
    p5 = tb5.text_frame.paragraphs[0]
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.add_run()
    r5.text = "→"
    set_font(r5, size=16, bold=True, color=BLUE)


# ============================================================
# 슬라이드 10: 해결 — 3계층 하네스 아키텍처
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "해결: 3계층 하네스 아키텍처 (Tool + Agent + Skill)"
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 3계층 다이어그램 (세로 스택)
LAYER_COLORS = [
    (RGBColor(0xE8, 0xEA, 0xF6), INDIGO, "Layer 3: Skills (오케스트레이터)",
     '사용자 요청 → 어떤 Flow를 실행할지 판단\n.claude/skills/*/skill.md'),
    (LIGHT_GREEN, GREEN, "Layer 2: Agents (지능형 판단)",
     'Flow 실행 중 예외 처리, 대안 탐색, 사용자 소통\n.claude/agents/*.md'),
    (RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0xFF, 0x8F, 0x00), "Layer 1: Tools & Flows (실행 레이어)  ← 신규",
     '실제 코드를 래핑한 실행 단위 + 체이닝 정의\ntools/*.yaml + flows/*.yaml + scripts/*.py'),
]

for i, (bg, border, title, desc) in enumerate(LAYER_COLORS):
    y = 1.4 + i * 1.5
    add_rect(slide, 0.5, y, 6.5, 1.3, bg, border)
    tb2 = add_textbox(slide, 0.7, y + 0.05, 6.1, 0.4)
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = title
    set_font(r2, size=14, bold=True, color=border)
    tb3 = add_textbox(slide, 0.7, y + 0.45, 6.1, 0.7)
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    r3 = p3.add_run()
    r3.text = desc
    set_font(r3, size=11, color=DARK_GRAY)

# 우측: Tool 정의 예시
add_rect(slide, 7.3, 1.4, 5.5, 2.6, WHITE, RGBColor(0xFF, 0x8F, 0x00))
tb = add_textbox(slide, 7.5, 1.5, 5.1, 0.35)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Tool 정의 예시: fetch_aia.yaml"
set_font(run, size=13, bold=True, color=RGBColor(0xFF, 0x8F, 0x00))

tb = add_textbox(slide, 7.5, 1.85, 5.1, 2.0)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, 'name: fetch_aia', size=10, color=DARK_GRAY, space_after=1)
add_paragraph(tf, 'type: python_script', size=10, color=GRAY, space_after=1)
add_paragraph(tf, 'script: scripts/fetch_aia.py', size=10, color=GRAY, space_after=1)
add_paragraph(tf, 'inputs:', size=10, color=DARK_GRAY, space_after=1)
add_paragraph(tf, '  - wavelength: int (94~335)', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '  - time_start, time_end: string', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '  - output_dir: path', size=10, color=GRAY, space_after=1)
add_paragraph(tf, 'outputs:', size=10, color=DARK_GRAY, space_after=1)
add_paragraph(tf, '  - fits_files: directory', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '  - manifest: json', size=10, color=GRAY, space_after=1)
add_paragraph(tf, 'dependencies: [drms, astropy]', size=10, color=GRAY)

# 우측 하단: Flow 정의 예시
add_rect(slide, 7.3, 4.2, 5.5, 2.6, WHITE, INDIGO)
tb = add_textbox(slide, 7.5, 4.3, 5.1, 0.35)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Flow 정의 예시: dem_pipeline.yaml"
set_font(run, size=13, bold=True, color=INDIGO)

tb = add_textbox(slide, 7.5, 4.65, 5.1, 2.0)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, 'name: dem_pipeline', size=10, color=DARK_GRAY, space_after=1)
add_paragraph(tf, 'steps:', size=10, color=DARK_GRAY, space_after=1)
add_paragraph(tf, '  1. fetch_data → tool: fetch_aia', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '  2. preprocess → tool: preprocess_aia', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '     input: {{steps.fetch_data.outputs.fits_files}}', size=10, color=RGBColor(0xFF, 0x8F, 0x00), space_after=1)
add_paragraph(tf, '  3. generate → tool: pix2pixcc_inference', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '  4. dem → tool: dem_inversion', size=10, color=GRAY, space_after=1)
add_paragraph(tf, '  5. visualize → tool: dem_visualizer', size=10, color=GRAY, space_after=3)
add_paragraph(tf, '→ Step간 데이터가 {{변수}}로 자동 연결', size=10, bold=True, color=INDIGO)

# 하단 좌측: 비교
add_rect(slide, 0.5, 5.9, 6.5, 0.9, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 0.7, 5.95, 6.1, 0.8)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Tool+Flow가 있으면: 검증된 코드 즉시 호출 · Step 자동 연결 · 재현성 보장 · 연구실 코드 직접 래핑"
set_font(run, size=13, bold=True, color=NAVY)


# ============================================================
# 슬라이드 11: 확장 가능성 — 100가지 하네스 아이디어
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "확장 가능성 — 100가지 특화 하네스"
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 카테고리별 표
headers = ["카테고리", "개수", "예시"]
data = [
    ["태양 데이터 처리 & 분석", "15", "AIA 파이프라인, synoptic map, 코로나홀 추적"],
    ["우주기상 예보 & 모니터링", "13", "Kp/Dst 예보, 태양풍 예측, SEP 예보"],
    ["딥러닝 모델 개발 & 실험", "15", "Pix2PixCC 학습, 초해상도, XAI 분석"],
    ["코로나/자기장 물리 연구", "10", "3D 복원, NLFFF, MHD, 코로나 진동"],
    ["논문 & 연구 지원", "13", "논문 초안, 문헌 검토, 제안서 작성"],
    ["데이터 관리 & 인프라", "10", "모델 레지스트리, 환경 복제, 재현성"],
    ["시각화 & 커뮤니케이션", "9", "태양 애니메이션, 3D 코로나, 인포그래픽"],
    ["교육 & 온보딩", "7", "신규 연구자 가이드, 튜토리얼 생성"],
    ["위성 미션 & 관측 지원", "5", "관측 계획, conjunction 탐색, L4 미션"],
    ["연구실 운영 & 관리", "3", "지식 베이스, 회의록, 로드맵"],
]
tbl = add_table(slide, 11, 3, 0.8, 1.4, 7.5, 4.8)
style_table(tbl, headers, data, col_widths=[2.5, 0.8, 4.2])

# 우측: 설명 박스
add_rect(slide, 8.8, 1.4, 4.0, 2.5, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 9.0, 1.5, 3.6, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "범용 vs 특화 하네스"
set_font(run, size=15, bold=True, color=NAVY)

tb = add_textbox(slide, 9.0, 1.9, 3.6, 1.8)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "현재 v2.0의 research-task는", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "범용 파이프라인으로 대부분을", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "처리할 수 있습니다.", size=11, color=DARK_GRAY, space_after=8)
add_paragraph(tf, "특화 하네스를 만들면:", size=11, bold=True, color=NAVY, space_after=4)
add_paragraph(tf, "- 정확한 코드 호출", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "- 도메인 지식 내장", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "- 오류율 감소", size=11, color=DARK_GRAY)

# 우측 하단: 핵심 메시지
add_rect(slide, 8.8, 4.2, 4.0, 2.0, RGBColor(0xFD, 0xFD, 0xE8), RGBColor(0xE0, 0xE0, 0xE0))
tb = add_textbox(slide, 9.0, 4.3, 3.6, 1.8)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, "하지만 이것은 샘플입니다", size=13, bold=True, color=NAVY, space_after=8)
add_paragraph(tf, "실제로 의미있는 하네스를 만들려면", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "연구실의 기존 코드를 읽고", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "그 코드에 맞는 하네스를", size=11, color=DARK_GRAY, space_after=2)
add_paragraph(tf, "자동 생성해야 합니다.", size=11, bold=True, color=NAVY)


# ============================================================
# 슬라이드 12: 코드 → 하네스 자동 생성 절차 (Tool/Flow 포함)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

AMBER = RGBColor(0xFF, 0x8F, 0x00)
LIGHT_AMBER = RGBColor(0xFF, 0xF8, 0xE1)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = '코드 헤리티지 → Tool/Flow + 하네스 자동 생성'
set_font(run, size=28, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 3.0, 0.05, INDIGO)

# 6개 Phase 박스 (가로 배열)
phases = [
    ("Phase 0", "코드 수집", "연구실 어디에\n뭐가 있는가", INDIGO),
    ("Phase 1", "코드 분석", "code-profiler가\nCode Profile 생성", INDIGO),
    ("Phase 2", "하네스 설계", "Code Profile →\nBlueprint 그루핑", INDIGO),
    ("Phase 2.5", "Tool/Flow 생성", "실행 가능한 도구\n+ 툴체인 YAML", AMBER),
    ("Phase 3", "에이전트 생성", "에이전트/스킬\n.md 파일 생성", INDIGO),
    ("Phase 4", "검증/배포", "드라이런 +\n사용자 피드백", INDIGO),
]

for i, (phase, title, desc, color) in enumerate(phases):
    x = 0.3 + i * 2.15
    add_rect(slide, x, 1.4, 1.95, 2.6, WHITE, color)
    add_rect(slide, x, 1.4, 1.95, 0.45, color)
    tb2 = add_textbox(slide, x, 1.4, 1.95, 0.45)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = phase
    set_font(r2, size=11, bold=True, color=WHITE)
    tb3 = add_textbox(slide, x + 0.05, 1.9, 1.85, 0.35)
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = title
    set_font(r3, size=12, bold=True, color=color)
    tb4 = add_textbox(slide, x + 0.05, 2.3, 1.85, 1.4)
    tb4.text_frame.word_wrap = True
    p4 = tb4.text_frame.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = desc
    set_font(r4, size=10, color=GRAY)

# 화살표
for i in range(5):
    x = 2.1 + i * 2.15
    tb5 = add_textbox(slide, x, 2.3, 0.4, 0.4)
    p5 = tb5.text_frame.paragraphs[0]
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.add_run()
    r5.text = "→"
    set_font(r5, size=16, bold=True, color=INDIGO)

# 하단 좌측: Phase 2.5 상세 (핵심 신규)
add_rect(slide, 0.3, 4.3, 6.3, 2.8, LIGHT_AMBER, RGBColor(0xFF, 0xCA, 0x28))
tb = add_textbox(slide, 0.5, 4.4, 5.9, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Phase 2.5: Tool/Flow 생성 (핵심 추가)"
set_font(run, size=15, bold=True, color=AMBER)

tb = add_textbox(slide, 0.5, 4.8, 5.9, 2.1)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, 'Code Profile의 각 실행 step:', size=11, bold=True, color=DARK_GRAY, space_after=2)
add_paragraph(tf, '  → Tool YAML (입출력 스키마 + 래퍼 스크립트)', size=11, color=GRAY, space_after=4)
add_paragraph(tf, 'Code Profile의 전체 파이프라인:', size=11, bold=True, color=DARK_GRAY, space_after=2)
add_paragraph(tf, '  → Flow YAML (Tool 체이닝 + {{변수}} 자동 연결)', size=11, color=GRAY, space_after=4)
add_paragraph(tf, '모델/가중치:', size=11, bold=True, color=DARK_GRAY, space_after=2)
add_paragraph(tf, '  → model_registry 등록 (model_card + tool.yaml)', size=11, color=GRAY, space_after=6)
add_paragraph(tf, '이 단계가 있어야 에이전트가 "설명"이 아닌 "실행"을 한다', size=11, bold=True, color=AMBER)

# 하단 우측: harness-factory
add_rect(slide, 6.8, 4.3, 6.0, 2.8, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 7.0, 4.4, 5.6, 0.4)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "harness-factory: 전체를 자동화하는 메타-하네스"
set_font(run, size=14, bold=True, color=NAVY)

tb = add_textbox(slide, 7.0, 4.8, 5.6, 2.1)
tf = tb.text_frame
tf.word_wrap = True
add_paragraph(tf, '4개 에이전트가 Phase 0~4를 자동 수행:', size=11, color=DARK_GRAY, space_after=6)
add_paragraph(tf, 'code-profiler — 코드 자동 분석', size=11, color=GRAY, space_after=2)
add_paragraph(tf, 'harness-designer — 하네스 구조 설계', size=11, color=GRAY, space_after=2)
add_paragraph(tf, 'tool-generator — Tool YAML + Flow YAML 생성', size=11, bold=True, color=AMBER, space_after=2)
add_paragraph(tf, 'harness-validator — 드라이런 검증 + 배포', size=11, color=GRAY, space_after=8)
add_paragraph(tf, '"코드 경로만 알려주면', size=12, bold=True, color=NAVY, space_after=1)
add_paragraph(tf, ' Tool + Flow + 에이전트가 자동으로 만들어집니다"', size=12, bold=True, color=NAVY)


# ============================================================
# 슬라이드 11: 납품 로드맵
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

tb = add_textbox(slide, 0.8, 0.4, 11, 0.7)
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = "납품 로드맵"
set_font(run, size=30, bold=True, color=NAVY)
add_rect(slide, 0.8, 1.05, 1.8, 0.05, INDIGO)

# 타임라인 박스들
roadmap = [
    ("Step 1", "코드 수집", "1~2일",
     "연구자 인터뷰 + 서버 스캔\ninventory.json 생성\n비버전관리 코드 식별"),
    ("Step 2", "코드 분석", "2~3일",
     "code-profiler로 Code Profile 생성\n입출력/의존성/실행방법 자동 추출\n코드당 10~30분"),
    ("Step 3", "하네스 설계", "1일",
     "Code Profile → Blueprint 그루핑\n에이전트/스킬 구조 설계\n연구자 리뷰 + 피드백"),
    ("Step 4", "자동 생성", "1일",
     "Blueprint → .md 파일 자동 생성\nmodel_card.md + run.sh\n코드당 5~10분"),
    ("Step 5", "검증/배포", "2~3일",
     "드라이런 (샘플 데이터)\n사용자 검증 + 피드백 반영\n최종 배포"),
]

for i, (step, title, duration, desc) in enumerate(roadmap):
    x = 0.5 + i * 2.55
    y = 1.4

    # 메인 박스
    add_rect(slide, x, y, 2.3, 3.8, WHITE, INDIGO)

    # Step 번호 바
    add_rect(slide, x, y, 2.3, 0.5, INDIGO)
    tb2 = add_textbox(slide, x, y, 2.3, 0.5)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = step
    set_font(r2, size=13, bold=True, color=WHITE)

    # 제목
    tb3 = add_textbox(slide, x + 0.1, y + 0.6, 2.1, 0.35)
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = title
    set_font(r3, size=14, bold=True, color=NAVY)

    # 기간 뱃지
    add_rect(slide, x + 0.6, y + 1.0, 1.1, 0.35, LIGHT_BG, TABLE_BORDER)
    tb4 = add_textbox(slide, x + 0.6, y + 1.0, 1.1, 0.35)
    p4 = tb4.text_frame.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = duration
    set_font(r4, size=11, bold=True, color=INDIGO)

    # 설명
    tb5 = add_textbox(slide, x + 0.15, y + 1.5, 2.0, 2.1)
    tb5.text_frame.word_wrap = True
    p5 = tb5.text_frame.paragraphs[0]
    p5.alignment = PP_ALIGN.LEFT
    r5 = p5.add_run()
    r5.text = desc
    set_font(r5, size=10, color=GRAY)

# 하단: 총 기간
box = add_rect(slide, 0.5, 5.5, 12.3, 0.8, LIGHT_BG, TABLE_BORDER)
tb = add_textbox(slide, 0.7, 5.55, 11.9, 0.7)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "총 예상 기간: 1~2주 (코드 규모에 따라 변동)"
set_font(run, size=18, bold=True, color=NAVY)

add_paragraph(tf, "현실: 연구실 코드의 80%는 README 없이 개인 PC에 존재 → Phase 1(코드 분석)이 핵심", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# 하단: 우선순위
tb = add_textbox(slide, 0.5, 6.5, 12.3, 0.6)
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "우선순위:  1순위 활발히 사용 중  →  2순위 논문 공개 코드  →  3순위 공통 유틸리티  →  4순위 졸업생 코드 보존"
set_font(run, size=12, color=DARK_GRAY)


# ============================================================
# 슬라이드 12: 기술 요약 + 마무리
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.333, 0.08, INDIGO)

# 기술 요약
tb = add_textbox(slide, 1, 0.5, 11.333, 0.7)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "기술 요약"
set_font(run, size=28, bold=True, color=WHITE)

# 스펙 박스들
specs = [
    ("기반 기술", "Claude Code\n(Anthropic AI)\n에이전트 프레임워크"),
    ("AI 모듈", "7개 전문 에이전트\n6개 스킬"),
    ("데이터", "NOAA, SDO, Solar Orbiter\nSTEREO, VSO"),
    ("모델 관리", "파일 기반 레지스트리\nmodel_card.md"),
    ("알림", "Slack Webhook\nTelegram Bot API"),
]

for i, (title, desc) in enumerate(specs):
    x = 0.8 + i * 2.5
    add_rect(slide, x, 1.5, 2.2, 2.2, INDIGO)
    tb2 = add_textbox(slide, x + 0.1, 1.6, 2.0, 0.4)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    set_font(r2, size=14, bold=True, color=RGBColor(0xC5, 0xCA, 0xE9))
    tb3 = add_textbox(slide, x + 0.1, 2.1, 2.0, 1.4)
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = desc
    set_font(r3, size=12, color=WHITE)

# 구분선
add_rect(slide, 5, 4.2, 3.333, 0.04, INDIGO)

# 메인 메시지
tb = add_textbox(slide, 1, 4.5, 11.333, 1.5)
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "SSWL AI Harness v2.0"
set_font(run, size=36, bold=True, color=WHITE)

add_paragraph(tf, "", size=8)
r = add_paragraph(tf, "연구실에서 개발된 모델을 누구나 활용하고,", size=16, color=RGBColor(0x9F, 0xA8, 0xDA), align=PP_ALIGN.CENTER, space_after=2)
add_paragraph(tf, "아이디어를 빠르게 검증합니다.", size=16, color=RGBColor(0x9F, 0xA8, 0xDA), align=PP_ALIGN.CENTER)

# 하단 안내문
tb = add_textbox(slide, 1, 6.5, 11.333, 0.5)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "본 시스템은 연구자의 판단을 대체하지 않습니다. AI는 도구이며, 최종 의사결정은 항상 연구자에게 있습니다."
set_font(run, size=11, color=RGBColor(0x79, 0x86, 0xCB), italic=True)


# ============================================================
# 저장
# ============================================================
output_path = "/home/youn_j/SSWL-harness-ops/PPT/pamphlet_r5.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
print(f"총 {len(prs.slides)} 슬라이드")
